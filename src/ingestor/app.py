from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
from pptx import Presentation
from pdfminer.high_level import extract_text as pdf_extract
import docx2txt
import tempfile
import os

app = FastAPI()

@app.post("/extract")
async def extract(file: UploadFile = File(...)):
    try:
        # Save uploaded file temporarily
        fd, tmp_path = tempfile.mkstemp(suffix=os.path.splitext(file.filename or "")[1])
        with os.fdopen(fd, "wb") as f:
            f.write(await file.read())

        filename = (file.filename or "").lower()
        full_text = ""
        outline = []

        # Extract based on file type
        if filename.endswith(".pptx"):
            prs = Presentation(tmp_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    slide_content = " ".join(slide_text)
                    full_text += slide_content + "\n\n"
                    # Add first 100 chars of each slide to outline
                    outline.append(f"Slide {slide_num}: {slide_content[:100]}")

        elif filename.endswith(".pdf"):
            full_text = pdf_extract(tmp_path)
            # Create outline from first line of each page (simplified)
            lines = [l.strip() for l in full_text.split("\n") if l.strip()]
            outline = lines[:50]  # First 50 non-empty lines as outline

        elif filename.endswith(".docx"):
            full_text = docx2txt.process(tmp_path)
            lines = [l.strip() for l in full_text.split("\n") if l.strip()]
            outline = lines[:50]

        else:
            os.remove(tmp_path)
            return JSONResponse(
                {"ok": False, "error": "Unsupported file type. Use .pptx, .pdf, or .docx"},
                status_code=400
            )

        os.remove(tmp_path)

        return JSONResponse({
            "ok": True,
            "filename": file.filename,
            "fullText": full_text.strip(),
            "outline": outline
        })

    except Exception as e:
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.remove(tmp_path)
        return JSONResponse({"ok": False, "error": str(e)}, status_code=500)

@app.get("/health")
async def health():
    return {"status": "healthy", "service": "ingestor"}
